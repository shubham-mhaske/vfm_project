#!/usr/bin/env python3
"""
Create VFM-Focused PowerPoint Presentation
Focus: Vision Foundation Models - SAM2, MedSAM, CLIP architectures and experiments

Usage:
    python scripts/analysis/create_detailed_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Project paths
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
figures_dir = os.path.join(project_root, 'results', 'figures')
output_dir = os.path.join(project_root, 'presentation_data')
os.makedirs(output_dir, exist_ok=True)

# Colors
TAMU_MAROON = RGBColor(80, 0, 0)
DARK_GRAY = RGBColor(51, 51, 51)
MEDIUM_GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 139, 34)
LIGHT_GREEN = RGBColor(230, 255, 230)
RED = RGBColor(180, 60, 60)
BLUE = RGBColor(30, 100, 180)
LIGHT_BLUE = RGBColor(230, 245, 255)
ORANGE = RGBColor(200, 120, 0)
LIGHT_ORANGE = RGBColor(255, 245, 230)


def add_title_bar(slide, title, subtitle=None):
    """Add maroon title bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAMU_MAROON
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.52), Inches(9.4), Inches(0.25))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(255, 210, 210)


def add_footer(slide, num, total):
    """Add footer and slide number."""
    footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(6), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Visual Foundation Models for Medical Image Segmentation • CSCE 689"
    p.font.size = Pt(9)
    p.font.color.rgb = MEDIUM_GRAY
    
    num_box = slide.shapes.add_textbox(Inches(9), Inches(7.05), Inches(0.7), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_text(slide, text, x, y, width, height, font_size=14, bold=False, color=None, align="left"):
    """Add text box."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color if color else DARK_GRAY
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return box


def add_bullets(slide, items, x, y, width, font_size=13):
    """Add bullet points."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("  "):
            p.text = f"   ◦ {item.strip()}"
            p.font.size = Pt(font_size - 1)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)
    return box


def add_table(slide, data, x, y, col_widths, highlight_row=None, font_size=11):
    """Add table."""
    rows, cols = len(data), len(data[0])
    total_w = sum(col_widths)
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), 
                                    Inches(total_w), Inches(0.35 * rows)).table
    
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)
    
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.alignment = PP_ALIGN.CENTER
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TAMU_MAROON
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                elif highlight_row and i == highlight_row:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREEN
                    para.font.bold = True
    return table


def add_image(slide, img_path, x, y, width=None, height=None):
    """Add image safely."""
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(width), height=Inches(height))
        elif width:
            slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(width))
        elif height:
            slide.shapes.add_picture(img_path, Inches(x), Inches(y), height=Inches(height))
        return True
    print(f"  Warning: {os.path.basename(img_path)} not found")
    return False


def add_model_box(slide, name, params, details, x, y, width, height, color):
    """Add a model description box."""
    bg = LIGHT_GREEN if color == GREEN else LIGHT_BLUE if color == BLUE else LIGHT_ORANGE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    add_text(slide, name, x + 0.15, y + 0.1, width - 0.3, 0.35, font_size=15, bold=True, color=color)
    add_text(slide, params, x + 0.15, y + 0.45, width - 0.3, 0.25, font_size=11, color=MEDIUM_GRAY)
    
    detail_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.75), Inches(width - 0.3), Inches(height - 0.9))
    tf = detail_box.text_frame
    tf.word_wrap = True
    for i, d in enumerate(details):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {d}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


# ============================================================================
# SLIDES
# ============================================================================

def slide_title(prs, num, total):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(2.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAMU_MAROON
    bar.line.fill.background()
    
    add_text(slide, "Visual Foundation Models for", 0.5, 2.75, 9, 0.6, font_size=28, color=WHITE, align="center")
    add_text(slide, "Medical Image Segmentation", 0.5, 3.35, 9, 0.7, font_size=38, bold=True, color=WHITE, align="center")
    add_text(slide, "SAM2  •  MedSAM  •  CLIP  •  Finetuning Experiments", 0.5, 5.1, 9, 0.4, font_size=16, color=DARK_GRAY, align="center")
    add_text(slide, "CSCE 689 • Fall 2025 • Texas A&M University", 0.5, 6.4, 9, 0.4, font_size=13, color=MEDIUM_GRAY, align="center")
    add_footer(slide, num, total)


def slide_vfm_overview(prs, num, total):
    """VFMs we explored."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Vision Foundation Models Explored", "Three VFM architectures for histopathology")
    add_footer(slide, num, total)
    
    # SAM2 box
    add_model_box(slide, "SAM2 (Segment Anything 2)", "224M parameters • Hiera-L backbone",
                  ["Hierarchical Vision Transformer", "Pretrained: SA-1B (11M images, 1B masks)", 
                   "Memory attention for video/image", "Promptable: points, boxes, masks"],
                  0.3, 1.1, 3.0, 2.4, BLUE)
    
    # MedSAM box  
    add_model_box(slide, "MedSAM", "93.7M parameters • ViT-B backbone",
                  ["Based on original SAM architecture", "Finetuned: 1.5M medical images",
                   "11 imaging modalities (CT, MRI, etc.)", "Box prompts optimized"],
                  3.5, 1.1, 3.0, 2.4, ORANGE)
    
    # CLIP box
    add_model_box(slide, "CLIP", "151M parameters • ViT-B/32",
                  ["Contrastive image-text learning", "Pretrained: 400M image-text pairs",
                   "Zero-shot transfer capability", "Text-based classification"],
                  6.7, 1.1, 3.0, 2.4, GREEN)
    
    # Architecture comparison table
    add_text(slide, "Architecture Comparison:", 0.3, 3.7, 9, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    data = [
        ["Model", "Backbone", "Params", "Pretraining", "Task"],
        ["SAM2", "Hiera-L", "224M", "SA-1B (natural)", "Segmentation"],
        ["MedSAM", "ViT-B", "93.7M", "Medical images", "Segmentation"],
        ["CLIP", "ViT-B/32", "151M", "Web image-text", "Classification"],
    ]
    add_table(slide, data, x=0.3, y=4.1, col_widths=[1.2, 1.3, 1.0, 2.0, 1.8])
    
    # Key insight
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5.9), Inches(9.4), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = MEDIUM_GRAY
    add_text(slide, "Research Question: How do general-purpose VFMs compare to domain-specific models on histopathology?", 
             0.5, 6.05, 9, 0.6, font_size=13, bold=True, color=DARK_GRAY, align="center")


def slide_sam2_architecture(prs, num, total):
    """SAM2 Architecture Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SAM2: Architecture & Approach", "Segment Anything Model 2 for histopathology")
    add_footer(slide, num, total)
    
    # Left: Architecture details
    add_text(slide, "Hiera-L Backbone:", 0.3, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Hierarchical Vision Transformer",
        "Multi-scale feature extraction",
        "224M total parameters",
        "  Image encoder: 212M",
        "  Mask decoder: 12M",
        "Efficient attention mechanism",
    ], x=0.3, y=1.35, width=4.5, font_size=12)
    
    add_text(slide, "Our Experimental Setup:", 0.3, 3.5, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Input: 1024×1024 histopathology tiles",
        "4 prompt strategies tested",
        "  Centroid points",
        "  Multi-point (5 boundary points)",
        "  Bounding boxes",
        "  Box + negative points",
        "Zero-shot evaluation (no training)",
    ], x=0.3, y=3.85, width=4.5, font_size=12)
    
    # Right: Qualitative figure
    add_image(slide, os.path.join(figures_dir, 'qualitative', 'qualitative_prompt_comparison.png'),
              5.0, 1.0, width=4.7, height=3.2)
    add_text(slide, "Visual comparison of input prompts and resulting segmentations", 
             5.0, 4.3, 4.7, 0.3, font_size=10, color=MEDIUM_GRAY, align="center")
    
    # Results summary
    add_text(slide, "SAM2 Results Summary:", 5.0, 4.7, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    data = [
        ["Input Type", "Dice Score"],
        ["Centroid", "0.338"],
        ["Multi-Point", "0.418"],
        ["Bounding Box", "0.553"],
        ["Box + Neg", "0.555 ✓"],
    ]
    add_table(slide, data, x=5.0, y=5.1, col_widths=[1.8, 1.2], highlight_row=4, font_size=10)


def slide_medsam_architecture(prs, num, total):
    """MedSAM Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "MedSAM: Domain-Specific VFM", "SAM finetuned on 1.5M medical images")
    add_footer(slide, num, total)
    
    # Left: Architecture
    add_text(slide, "Model Architecture:", 0.3, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Based on SAM (ViT-B backbone)",
        "93.7M parameters (2.4× smaller than SAM2)",
        "Finetuned on diverse medical data:",
        "  CT, MRI, Ultrasound, X-ray",
        "  Dermoscopy, Endoscopy, Pathology",
        "  1.5M images, 11 modalities",
        "Optimized for bounding box input",
    ], x=0.3, y=1.35, width=4.5, font_size=12)
    
    add_text(slide, "Test-Time Augmentation (TTA):", 0.3, 3.7, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "4 geometric transforms:",
        "  Original, H-flip, V-flip, HV-flip",
        "Average predictions for robustness",
        "Improves Dice: 0.522 → 0.536 (+2.7%)",
    ], x=0.3, y=4.05, width=4.5, font_size=12)
    
    # Right: Comparison
    add_text(slide, "SAM2 vs MedSAM Head-to-Head:", 5.0, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    data = [
        ["Aspect", "SAM2", "MedSAM"],
        ["Parameters", "224M", "93.7M"],
        ["Backbone", "Hiera-L", "ViT-B"],
        ["Pretraining", "Natural", "Medical"],
        ["Dice Score", "0.555", "0.536"],
        ["IoU Score", "0.408", "0.389"],
    ]
    add_table(slide, data, x=5.0, y=1.4, col_widths=[1.5, 1.3, 1.3], highlight_row=4)
    
    # Key finding box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.8), Inches(4.5), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.color.rgb = GREEN
    add_text(slide, "Key Finding:", 5.2, 3.95, 4.1, 0.3, font_size=12, bold=True, color=GREEN)
    add_text(slide, "SAM2 (general-purpose, 2.4× larger) outperforms MedSAM (domain-specific) by 3.5% on histopathology. Model capacity may matter more than domain-specific pretraining.", 
             5.2, 4.25, 4.1, 0.9, font_size=11, color=DARK_GRAY)
    
    # Qualitative
    add_image(slide, os.path.join(figures_dir, 'qualitative', 'qualitative_method_comparison.png'),
              0.3, 5.3, width=9.4, height=1.6)


def slide_finetuning_experiments(prs, num, total):
    """Finetuning experiments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "SAM2 Finetuning Experiments", "Does domain adaptation improve performance?")
    add_footer(slide, num, total)
    
    # Left: Training approaches
    add_text(slide, "Training Configurations:", 0.3, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    # Full finetuning box
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(4.5), Inches(1.3))
    b1.fill.solid()
    b1.fill.fore_color.rgb = LIGHT_BLUE
    b1.line.color.rgb = BLUE
    add_text(slide, "Full Finetuning", 0.5, 1.5, 4, 0.3, font_size=13, bold=True, color=BLUE)
    add_bullets(slide, [
        "All 224M parameters updated",
        "BCE Loss: 100 epochs",
        "Focal Loss: 50 epochs (γ=2)",
    ], x=0.5, y=1.85, width=4, font_size=11)
    
    # LoRA box
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(2.8), Inches(4.5), Inches(1.3))
    b2.fill.solid()
    b2.fill.fore_color.rgb = LIGHT_ORANGE
    b2.line.color.rgb = ORANGE
    add_text(slide, "LoRA Adaptation", 0.5, 2.9, 4, 0.3, font_size=13, bold=True, color=ORANGE)
    add_bullets(slide, [
        "Low-rank adapters (r=8)",
        "Only 0.5% params trainable",
        "30 epochs training",
    ], x=0.5, y=3.25, width=4, font_size=11)
    
    # Dataset info
    add_text(slide, "Dataset: BCSS", 0.3, 4.3, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "85 training / 21 validation / 45 test",
        "5 tissue classes",
        "1024×1024 patches",
    ], x=0.3, y=4.65, width=4.5, font_size=11)
    
    # Right: Results
    add_text(slide, "Test Set Results:", 5.0, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    data = [
        ["Method", "Dice", "IoU", "Δ vs Zero-Shot"],
        ["Zero-Shot", "0.555", "0.408", "—"],
        ["BCE Finetuned", "0.371", "0.264", "-33%"],
        ["Focal Finetuned", "0.372", "0.265", "-33%"],
        ["LoRA (r=8)", "0.355", "0.248", "-36%"],
    ]
    add_table(slide, data, x=5.0, y=1.4, col_widths=[1.6, 0.7, 0.7, 1.2], highlight_row=1)
    
    # Analysis box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.5), Inches(4.5), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.color.rgb = GREEN
    add_text(slide, "Analysis: Why Zero-Shot Wins", 5.2, 3.65, 4.1, 0.3, font_size=12, bold=True, color=GREEN)
    add_bullets(slide, [
        "Small dataset (n=85) causes overfitting",
        "SA-1B pretrained features are robust",
        "Catastrophic forgetting of general knowledge",
        "VFM generalization > domain-specific training",
    ], x=5.2, y=4.0, width=4, font_size=11)
    
    # Training curves
    add_image(slide, os.path.join(figures_dir, 'academic', 'fig3_training_analysis.png'),
              0.3, 5.6, width=9.4, height=1.35)


def slide_clip_classification(prs, num, total):
    """CLIP for classification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "CLIP: Zero-Shot Classification", "Contrastive Language-Image Pretraining")
    add_footer(slide, num, total)
    
    # Left: Architecture
    add_text(slide, "CLIP Architecture:", 0.3, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Dual-encoder architecture:",
        "  Vision: ViT-B/32 (86M params)",
        "  Text: Transformer (63M params)",
        "Contrastive learning objective",
        "Pretrained on 400M image-text pairs",
        "Zero-shot via text similarity",
    ], x=0.3, y=1.35, width=4.5, font_size=12)
    
    add_text(slide, "Our Classification Pipeline:", 0.3, 3.2, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "1. SAM2 segments tissue regions",
        "2. Crop regions from original image",
        "3. CLIP encodes cropped regions",
        "4. Match against class text embeddings",
        "5. Assign class with highest similarity",
    ], x=0.3, y=3.55, width=4.5, font_size=12)
    
    # Right: Text encoding strategies
    add_text(slide, "Text Encoding Approaches:", 5.0, 1.0, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    # Manual box
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.4), Inches(4.5), Inches(1.0))
    b1.fill.solid()
    b1.fill.fore_color.rgb = LIGHT_GRAY
    b1.line.color.rgb = MEDIUM_GRAY
    add_text(slide, "Manual Descriptions", 5.2, 1.5, 4, 0.25, font_size=12, bold=True, color=DARK_GRAY)
    add_text(slide, "Hand-crafted text for each tissue class", 5.2, 1.8, 4, 0.25, font_size=10, color=MEDIUM_GRAY)
    
    # LLM box
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.5), Inches(4.5), Inches(1.0))
    b2.fill.solid()
    b2.fill.fore_color.rgb = LIGHT_GREEN
    b2.line.color.rgb = GREEN
    add_text(slide, "LLM-Generated (GPT-4/Gemini)", 5.2, 2.6, 4, 0.25, font_size=12, bold=True, color=GREEN)
    add_text(slide, "Automated description generation at scale", 5.2, 2.9, 4, 0.25, font_size=10, color=DARK_GRAY)
    
    # Results
    add_text(slide, "Classification Results:", 5.0, 3.7, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    data = [
        ["Text Source", "Accuracy"],
        ["LLM Few-Shot", "44.4%"],
        ["Manual v2", "42.2%"],
        ["LLM Jargon", "12.2%"],
        ["Random Baseline", "20.0%"],
    ]
    add_table(slide, data, x=5.0, y=4.05, col_widths=[2.0, 1.2], highlight_row=1, font_size=11)
    
    # CLIP figure
    add_image(slide, os.path.join(figures_dir, 'academic', 'fig2_clip_analysis.png'),
              0.3, 5.6, width=9.4, height=1.35)


def slide_clip_confusion_matrix(prs, num, total):
    """CLIP Confusion Matrix Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "CLIP Classification: Error Analysis", "Confusion Matrix reveals class confusions")
    add_footer(slide, num, total)
    
    # Confusion matrix image
    add_image(slide, os.path.join(figures_dir, 'academic', 'clip_confusion_matrix.png'),
              0.3, 1.0, height=4.0)
    
    # Key findings
    add_text(slide, "Key Observations:", 5.0, 1.0, 4.7, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Tumor: 82% recall (best class)",
        "Blood Vessel: 100% recall",
        "  but 75% false positive rate!",
        "Lymphocyte: 0% recall (CLIP fails)",
        "Necrosis: 0% recall (subtle texture)",
        "",
        "Why this pattern?",
        "• Blood vessels have distinct shape",
        "• Lymphocytes lack visual structure",
        "• CLIP prefers structural over texture",
    ], x=5.0, y=1.35, width=4.5, font_size=11)
    
    # Analysis box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(5.1), Inches(4.5), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_ORANGE
    box.line.color.rgb = ORANGE
    add_text(slide, "Implication for Medical AI:", 5.2, 5.25, 4.1, 0.3, font_size=12, bold=True, color=ORANGE)
    add_text(slide, "CLIP's web-scale pretraining creates bias toward structural features. Fine-grained texture discrimination requires domain-specific training or better prompts.", 
             5.2, 5.55, 4.1, 0.7, font_size=10, color=DARK_GRAY)


def slide_prompt_evolution(prs, num, total):
    """Prompt Evolution slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Prompt Engineering: The Key to CLIP", "From medical jargon to visual descriptors")
    add_footer(slide, num, total)
    
    # Prompt evolution image
    add_image(slide, os.path.join(figures_dir, 'academic', 'prompt_evolution.png'),
              0.2, 1.0, width=9.6, height=5.0)
    
    # Bottom insight
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.1), Inches(9.4), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.color.rgb = GREEN
    add_text(slide, "Key Insight: CLIP understands visual language, not medical terminology. Describing colors and shapes (\"dark purple dots\") works better than pathology terms (\"hyperchromatic nuclei\").", 
             0.5, 6.25, 9, 0.5, font_size=12, bold=True, color=DARK_GRAY, align="center")


def slide_failed_experiments(prs, num, total):
    """Failed Experiments Table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Finetuning: Systematic Failure Analysis", "All domain adaptation attempts underperformed zero-shot")
    add_footer(slide, num, total)
    
    # Failed experiments image
    add_image(slide, os.path.join(figures_dir, 'academic', 'failed_experiments.png'),
              0.2, 1.0, width=9.6, height=4.5)
    
    # Key insight
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5.6), Inches(9.4), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = BLUE
    add_text(slide, "Why Zero-Shot Wins:", 0.5, 5.75, 4.5, 0.3, font_size=12, bold=True, color=BLUE)
    add_bullets(slide, [
        "SA-1B features are robust and general",
        "Small dataset (n=85) causes overfitting",
        "Catastrophic forgetting destroys pretrained knowledge",
    ], x=0.5, y=6.05, width=4.3, font_size=11)
    
    add_text(slide, "Recommendation:", 5.0, 5.75, 4.5, 0.3, font_size=12, bold=True, color=BLUE)
    add_bullets(slide, [
        "Use zero-shot VFMs on small datasets",
        "Finetuning needs 1000+ diverse images",
        "Consider parameter-efficient methods (LoRA) with regularization",
    ], x=5.0, y=6.05, width=4.3, font_size=11)


def slide_perclass_analysis(prs, num, total):
    """Per-class performance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Per-Class Analysis: Model Behavior", "How VFMs perform on different tissue types")
    add_footer(slide, num, total)
    
    # Left: Qualitative
    add_image(slide, os.path.join(figures_dir, 'qualitative', 'qualitative_per_class.png'),
              0.2, 1.0, width=5.3, height=4.2)
    
    # Right: Segmentation results
    add_text(slide, "SAM2 Segmentation by Class:", 5.7, 1.0, 4, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    data = [
        ["Tissue", "Dice", "IoU", "N"],
        ["Necrosis", "0.699", "0.567", "23"],
        ["Tumor", "0.560", "0.419", "45"],
        ["Lymphocyte", "0.549", "0.408", "37"],
        ["Stroma", "0.537", "0.391", "45"],
        ["Blood Vessel", "0.504", "0.362", "31"],
    ]
    add_table(slide, data, x=5.7, y=1.4, col_widths=[1.3, 0.65, 0.65, 0.5], highlight_row=1, font_size=10)
    
    # CLIP results
    add_text(slide, "CLIP Classification by Class:", 5.7, 3.5, 4, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    
    data2 = [
        ["Tissue", "Precision", "Recall", "F1"],
        ["Tumor", "0.65", "0.72", "0.68"],
        ["Stroma", "0.45", "0.55", "0.49"],
        ["Necrosis", "0.52", "0.48", "0.50"],
        ["Lymphocyte", "0.18", "0.12", "0.14"],
        ["Blood Vessel", "0.33", "0.25", "0.28"],
    ]
    add_table(slide, data2, x=5.7, y=3.85, col_widths=[1.3, 0.8, 0.7, 0.5], highlight_row=1, font_size=10)
    
    # Analysis
    add_text(slide, "Observations:", 5.7, 5.8, 4, 0.35, font_size=12, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Distinct textures (necrosis) → best",
        "Small structures (vessels) → hardest",
        "CLIP struggles with lymphocytes",
    ], x=5.7, y=6.15, width=3.8, font_size=10)


def slide_qualitative_results(prs, num, total):
    """Full qualitative comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Qualitative Results: Method Comparison")
    add_footer(slide, num, total)
    
    # Main figure
    add_image(slide, os.path.join(figures_dir, 'qualitative', 'qualitative_method_comparison.png'),
              0.15, 1.0, width=9.7, height=5.2)
    
    add_text(slide, "Left→Right: Original | Ground Truth | SAM2 (Centroid) | SAM2 (Box) | SAM2 (Box+Neg) | MedSAM", 
             0.3, 6.35, 9.4, 0.35, font_size=11, color=MEDIUM_GRAY, align="center")


def slide_findings(prs, num, total):
    """Key findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Findings: VFM Insights")
    add_footer(slide, num, total)
    
    findings = [
        ("1", "Model Capacity Matters", "SAM2 224M > MedSAM 93.7M", 
         "Larger general VFM outperforms smaller domain-specific model by 3.5%", BLUE),
        ("2", "Zero-Shot > Finetuning", "0.555 vs 0.372 Dice", 
         "On small datasets, pretrained features preserve better than finetuning", GREEN),
        ("3", "Input Representation Critical", "+64% with boxes", 
         "How we encode spatial information dramatically affects VFM performance", ORANGE),
        ("4", "CLIP Requires Good Alignment", "44.4% accuracy", 
         "Text descriptions must match CLIP's visual vocabulary for zero-shot transfer", GREEN),
        ("5", "Negative Transfer Observed", "-33% after finetuning", 
         "Domain adaptation can hurt when training data is limited", RED),
    ]
    
    y = 1.1
    for num_id, title, metric, detail, color in findings:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(y), Inches(0.45), Inches(0.45))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()
        add_text(slide, num_id, 0.3, y + 0.08, 0.45, 0.35, font_size=16, bold=True, color=WHITE, align="center")
        add_text(slide, title, 0.9, y, 5, 0.35, font_size=14, bold=True, color=DARK_GRAY)
        add_text(slide, metric, 7.5, y, 2.2, 0.35, font_size=14, bold=True, color=color, align="right")
        add_text(slide, detail, 0.9, y + 0.38, 8.8, 0.3, font_size=11, color=MEDIUM_GRAY)
        y += 1.05


def slide_conclusion(prs, num, total):
    """Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Conclusion: VFM Recommendations")
    add_footer(slide, num, total)
    
    # Best config
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.4), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.color.rgb = GREEN
    box.line.width = Pt(3)
    
    add_text(slide, "Recommended VFM Configuration for Histopathology", 0.5, 1.2, 9, 0.4, font_size=16, bold=True, color=GREEN)
    add_bullets(slide, [
        "SAM2 (Hiera-L) for segmentation → 0.555 Dice, best across all experiments",
        "CLIP (ViT-B/32) for classification → 44.4% accuracy, 2× random baseline",
        "Zero-shot inference → No finetuning needed, use pretrained weights directly",
    ], x=0.5, y=1.7, width=9, font_size=12)
    
    # Implications
    add_text(slide, "Implications for VFM Research:", 0.3, 2.9, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "General VFMs transfer well to medical",
        "Larger models can beat specialized ones",
        "Small datasets favor zero-shot",
        "Input encoding matters significantly",
    ], x=0.3, y=3.25, width=4.5, font_size=12)
    
    # Future work
    add_text(slide, "Future Directions:", 5.2, 2.9, 4.5, 0.35, font_size=14, bold=True, color=TAMU_MAROON)
    add_bullets(slide, [
        "Larger datasets for finetuning study",
        "BiomedCLIP, PathSAM comparison",
        "Multi-task VFM architectures",
        "Uncertainty estimation in VFMs",
    ], x=5.2, y=3.25, width=4.5, font_size=12)
    
    # Summary figure
    add_image(slide, os.path.join(figures_dir, 'academic', 'fig5_summary_results.png'),
              0.3, 5.1, width=9.4, height=1.85)


def slide_thank_you(prs, num, total):
    """Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.4), Inches(10), Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAMU_MAROON
    bar.line.fill.background()
    
    add_text(slide, "Thank You!", 0.5, 2.8, 9, 0.8, font_size=44, bold=True, color=WHITE, align="center")
    add_text(slide, "Questions?", 0.5, 3.7, 9, 0.5, font_size=26, color=RGBColor(255, 200, 200), align="center")
    
    add_text(slide, "SAM2 • MedSAM • CLIP • Zero-Shot > Finetuning • 17 Experiments", 
             0.5, 5.2, 9, 0.4, font_size=14, color=MEDIUM_GRAY, align="center")
    
    add_footer(slide, num, total)


# ============================================================================
# MAIN
# ============================================================================

def create_full_presentation():
    """Create VFM-focused presentation."""
    print("=" * 60)
    print("Creating VFM-Focused Presentation")
    print("=" * 60)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides = [
        ("Title", slide_title),
        ("VFMs Overview", slide_vfm_overview),
        ("SAM2 Architecture", slide_sam2_architecture),
        ("MedSAM Architecture", slide_medsam_architecture),
        ("Finetuning Experiments", slide_finetuning_experiments),
        ("Failed Experiments", slide_failed_experiments),
        ("CLIP Classification", slide_clip_classification),
        ("CLIP Confusion Matrix", slide_clip_confusion_matrix),
        ("Prompt Evolution", slide_prompt_evolution),
        ("Per-Class Analysis", slide_perclass_analysis),
        ("Qualitative Results", slide_qualitative_results),
        ("Key Findings", slide_findings),
        ("Conclusion", slide_conclusion),
        ("Thank You", slide_thank_you),
    ]
    
    total = len(slides)
    for i, (name, func) in enumerate(slides, 1):
        print(f"  Slide {i}/{total}: {name}")
        func(prs, i, total)
    
    output_path = os.path.join(output_dir, 'VFM_Detailed_Presentation.pptx')
    prs.save(output_path)
    
    print("=" * 60)
    print(f"✅ Saved: {output_path}")
    print("=" * 60)
    return output_path


if __name__ == '__main__':
    create_full_presentation()
