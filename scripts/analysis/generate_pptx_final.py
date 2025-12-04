#!/usr/bin/env python3
"""
Generate Final Presentation: Promptable Pathology

Uses ONLY existing figures from results/figures/academic/ and results/figures/qualitative/
No new figure generation - just PowerPoint creation.

Usage:
    python scripts/analysis/generate_pptx_final.py
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Project paths
PROJECT_ROOT = Path(__file__).parent.parent.parent
ACADEMIC_FIGS = PROJECT_ROOT / 'results' / 'figures' / 'academic'
QUALITATIVE_FIGS = PROJECT_ROOT / 'results' / 'figures' / 'qualitative'
PRESENTATION_FIGS = PROJECT_ROOT / 'results' / 'figures' / 'presentation'
OUTPUT_DIR = PROJECT_ROOT / 'presentation_data'

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def find_image(base_path, filename):
    """Find image file, preferring .png over .pdf."""
    png_path = base_path / filename
    if png_path.exists():
        return str(png_path)
    
    # Try without extension
    stem = Path(filename).stem
    for ext in ['.png', '.pdf', '.jpg', '.jpeg']:
        candidate = base_path / f"{stem}{ext}"
        if candidate.exists():
            return str(candidate)
    
    # Check alternative directories
    for alt_dir in [ACADEMIC_FIGS, QUALITATIVE_FIGS, PRESENTATION_FIGS]:
        for ext in ['.png', '.pdf']:
            candidate = alt_dir / f"{stem}{ext}"
            if candidate.exists():
                return str(candidate)
    
    print(f"  Warning: Image not found: {filename}")
    return None


def create_presentation():
    """Create the 8-slide final presentation."""
    print("=" * 70)
    print("Generating: Promptable_Pathology_Final.pptx")
    print("=" * 70)
    
    # Create widescreen presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide dimensions
    SLIDE_WIDTH = 13.333
    SLIDE_HEIGHT = 7.5
    
    # Color scheme
    DARK_BLUE = RGBColor(44, 62, 80)
    GREEN = RGBColor(39, 174, 96)
    GRAY = RGBColor(127, 140, 141)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(236, 240, 241)
    
    def add_title_bar(slide, title_text, color=DARK_BLUE):
        """Add a colored title bar at the top of the slide."""
        # Title bar background
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(SLIDE_WIDTH), Inches(0.9)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15),
            Inches(SLIDE_WIDTH - 0.8), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        return 1.0  # Return content start position
    
    def add_image_centered(slide, image_path, top, max_width=12.5, max_height=4.5):
        """Add an image centered on the slide."""
        if not image_path or not os.path.exists(image_path):
            print(f"    Skipping missing image: {image_path}")
            return
        
        try:
            # Add picture and let pptx calculate size, then we'll adjust
            pic = slide.shapes.add_picture(
                image_path,
                Inches(0.4), Inches(top),
                width=Inches(max_width)
            )
            
            # Center horizontally
            pic.left = int((Inches(SLIDE_WIDTH) - pic.width) / 2)
            
            print(f"    Added: {Path(image_path).name}")
        except Exception as e:
            print(f"    Error adding {image_path}: {e}")
    
    def add_key_points(slide, points, top, left=0.5):
        """Add bullet points to the slide."""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(SLIDE_WIDTH - 1), Inches(2.0)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE
            p.space_after = Pt(8)
    
    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    print("\n[1/8] Title Slide...")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background gradient effect (solid for simplicity)
    bg = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 252)
    bg.line.fill.background()
    
    # Main title
    title_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(2.2),
        Inches(SLIDE_WIDTH - 1), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Promptable Pathology"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle line 1
    sub1 = slide1.shapes.add_textbox(
        Inches(0.5), Inches(3.5),
        Inches(SLIDE_WIDTH - 1), Inches(0.6)
    )
    tf = sub1.text_frame
    p = tf.paragraphs[0]
    p.text = "Zero-Shot Medical Image Segmentation"
    p.font.size = Pt(28)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle line 2 (Course info)
    sub2 = slide1.shapes.add_textbox(
        Inches(0.5), Inches(4.3),
        Inches(SLIDE_WIDTH - 1), Inches(0.5)
    )
    tf = sub2.text_frame
    p = tf.paragraphs[0]
    p.text = "CSCE 689 Fall 2025 | Shubham Mhaske"
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # SLIDE 2: Methodology
    # =========================================================================
    print("[2/8] Methodology Slide...")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide2, "Two-Stage Pipeline: SAM2 + CLIP")
    
    img_path = find_image(ACADEMIC_FIGS, 'fig4_method_overview.png')
    add_image_centered(slide2, img_path, content_top + 0.2, max_width=12, max_height=4.5)
    
    add_key_points(slide2, [
        "Modular Two-Stage Design: SAM2 (Segmentation) + CLIP (Classification)",
        "User provides prompts (points/boxes) → System segments and classifies tissue regions",
        "Dataset: BCSS (151 images, 5 tissue classes: tumor, stroma, lymphocyte, necrosis, blood vessel)"
    ], top=5.5)
    
    # =========================================================================
    # SLIDE 3: Failed Experiments (Iceberg)
    # =========================================================================
    print("[3/8] Failed Experiments Slide...")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide3, "Why Finetuning Failed")
    
    img_path = find_image(ACADEMIC_FIGS, 'fig3_training_analysis.png')
    add_image_centered(slide3, img_path, content_top + 0.1, max_width=11, max_height=4.0)
    
    add_key_points(slide3, [
        "Tested: PathSAM2+CTransPath (0.37), LoRA adapters (0.27-0.36), Focal Loss (0.37)",
        "Linear Probe: Logistic Regression on CLIP features (40.4%) — still < Zero-Shot (44.4%)",
        "Zero-Shot SAM2 (0.55) outperforms ALL finetuned models by 30-50%"
    ], top=5.2)
    
    # =========================================================================
    # SLIDE 4: Segmentation Quantitative
    # =========================================================================
    print("[4/8] Segmentation Quantitative Slide...")
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide4, "Prompt Strategy Comparison")
    
    img_path = find_image(ACADEMIC_FIGS, 'fig1_segmentation_comprehensive.png')
    add_image_centered(slide4, img_path, content_top + 0.1, max_width=12, max_height=4.2)
    
    add_key_points(slide4, [
        "Box prompts: 0.55 Dice vs Centroid: 0.34 (+64% improvement)",
        "TTA with Prompt Rotation: Geometrically transforming box coords to match augmentations (+2%)",
        "Ours (0.555) beats fully supervised MedSAM ViT-B Model (0.536)"
    ], top=5.4)
    
    # =========================================================================
    # SLIDE 5: Visual Evidence (Qualitative)
    # =========================================================================
    print("[5/8] Visual Evidence Slide...")
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide5, "Visual Comparison Across Methods")
    
    # Try presentation figures first, then qualitative
    img_path = find_image(PRESENTATION_FIGS, 'best_method_comparison.png')
    if not img_path:
        img_path = find_image(QUALITATIVE_FIGS, 'qualitative_method_comparison.png')
    add_image_centered(slide5, img_path, content_top + 0.1, max_width=12.5, max_height=4.5)
    
    add_key_points(slide5, [
        "Best samples: Ours (Box+Neg) achieves tight boundaries vs blob-like point prompts",
        "Centroid prompts fail on multi-region classes; Box captures full extent",
        "Negative points suppress background bleeding in complex tissue interfaces"
    ], top=5.7)
    
    # =========================================================================
    # SLIDE 6: Prompt Engineering Evolution
    # =========================================================================
    print("[6/8] Prompt Engineering Slide...")
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide6, "CLIP Prompt Engineering")
    
    img_path = find_image(ACADEMIC_FIGS, 'fig2_clip_analysis.png')
    add_image_centered(slide6, img_path, content_top + 0.1, max_width=11.5, max_height=4.2)
    
    add_key_points(slide6, [
        "V1 'Medical Jargon' (12%) → V3 'Gemini Few-Shot' (44%) = +264% improvement",
        "Key: Translate pathology terms to visual descriptors CLIP understands",
        "Also tested: PLIP (27%), Ensembles (31%) - CLIP alone is best"
    ], top=5.4)
    
    # =========================================================================
    # SLIDE 7: Per-Class Analysis
    # =========================================================================
    print("[7/8] Per-Class Analysis Slide...")
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide7, "Class Imbalance & Performance")
    
    # Try best_per_class first, then qualitative_per_class
    img_path = find_image(PRESENTATION_FIGS, 'best_per_class.png')
    if not img_path:
        img_path = find_image(QUALITATIVE_FIGS, 'qualitative_per_class.png')
    add_image_centered(slide7, img_path, content_top + 0.1, max_width=12, max_height=4.3)
    
    add_key_points(slide7, [
        "Best: Tumor (0.70 Dice), Stroma (0.60) | Worst: Lymphocytes (0.25), Blood (0.42)",
        "Class imbalance (Stroma=54%, Lymphocytes=3%) drives performance gap",
        "Rare classes need domain-specific prompts or targeted augmentation"
    ], top=5.5)
    
    # =========================================================================
    # SLIDE 8: Summary & Key Takeaways
    # =========================================================================
    print("[8/8] Summary Slide...")
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_bar(slide8, "Key Findings & Takeaways", color=GREEN)
    
    img_path = find_image(ACADEMIC_FIGS, 'fig5_summary_results.png')
    if img_path:
        add_image_centered(slide8, img_path, content_top + 0.1, max_width=11, max_height=3.5)
        points_top = 4.8
    else:
        points_top = 1.5
    
    # Key takeaways with emphasis
    takeaways = [
        ("Key Finding", "For small medical datasets (N<100), Prompt Engineering > Model Finetuning"),
        ("Segmentation", "Zero-Shot SAM2 + Box+Neg: 0.555 Dice (beats finetuned MedSAM 0.536)"),
        ("Classification", "CLIP + Gemini Few-Shot Prompts: 44.4% Accuracy (+264% vs jargon)"),
        ("Insight", "50+ failed experiments → Simple zero-shot with good prompts wins"),
    ]
    
    y_pos = points_top
    for label, text in takeaways:
        # Label
        label_box = slide8.shapes.add_textbox(
            Inches(0.5), Inches(y_pos),
            Inches(2.0), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ {label}:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        # Text
        text_box = slide8.shapes.add_textbox(
            Inches(2.5), Inches(y_pos),
            Inches(10.5), Inches(0.4)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        
        y_pos += 0.55
    
    # =========================================================================
    # Save Presentation
    # =========================================================================
    output_path = OUTPUT_DIR / 'Promptable_Pathology_Final.pptx'
    prs.save(str(output_path))
    
    print("\n" + "=" * 70)
    print(f"SUCCESS! Saved: {output_path}")
    print("=" * 70)
    
    return str(output_path)


if __name__ == '__main__':
    create_presentation()
