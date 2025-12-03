#!/usr/bin/env python3
"""
Create PowerPoint Presentation for VFM Project
Generates a 12-slide academic presentation with figures and results.

Usage:
    pip install python-pptx
    python scripts/analysis/create_presentation.py
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Project paths
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
figures_dir = os.path.join(project_root, 'results', 'figures')
output_dir = os.path.join(project_root, 'presentation_data')
os.makedirs(output_dir, exist_ok=True)

# Colors
TAMU_MAROON = RGBColor(80, 0, 0)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(178, 34, 34)


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TAMU_MAROON
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CSCE 689 • Fall 2025 • Texas A&M University"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullets=None, image_path=None, image_position='right'):
    """Add a content slide with optional bullets and image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TAMU_MAROON
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    if image_path and os.path.exists(image_path):
        if bullets:
            # Left: bullets, Right: image
            if image_position == 'right':
                bullet_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.2), Inches(5.5))
                slide.shapes.add_picture(image_path, Inches(4.7), Inches(1.2), width=Inches(5.1))
            else:  # image on left
                slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.2), width=Inches(5.1))
                bullet_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4.2), Inches(5.5))
            
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(12)
        else:
            # Full-width image
            slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.2), width=Inches(9.4))
    elif bullets:
        # Full-width bullets
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(16)
    
    return slide


def add_figure_slide(prs, title, image_path, caption=None):
    """Add a slide with a large figure."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TAMU_MAROON
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Image
    if os.path.exists(image_path):
        # Center the image
        slide.shapes.add_picture(image_path, Inches(0.2), Inches(0.9), width=Inches(9.6))
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_results_slide(prs, title, results_table, highlight_row=None):
    """Add a slide with a results table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TAMU_MAROON
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    rows = len(results_table)
    cols = len(results_table[0]) if results_table else 0
    
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5 * rows)).table
    
    for i, row_data in enumerate(results_table):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            
            # Header row styling
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TAMU_MAROON
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.size = Pt(14)
            else:
                # Highlight best row
                if highlight_row and i == highlight_row:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 255, 220)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13)
                    paragraph.font.color.rgb = DARK_GRAY
    
    return slide


def add_key_findings_slide(prs, findings):
    """Add a key findings slide with icons."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TAMU_MAROON
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Findings boxes
    y_pos = 1.2
    for i, (finding, value, context) in enumerate(findings):
        # Number box
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.6), Inches(0.6))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = TAMU_MAROON
        num_shape.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(0.6), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Finding text
        finding_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(5.5), Inches(0.6))
        tf = finding_box.text_frame
        p = tf.paragraphs[0]
        p.text = finding
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(7), Inches(y_pos), Inches(2.5), Inches(0.4))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.RIGHT
        
        # Context
        ctx_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.5), Inches(8), Inches(0.4))
        tf = ctx_box.text_frame
        p = tf.paragraphs[0]
        p.text = context
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        
        y_pos += 1.4
    
    return slide


def create_presentation():
    """Create the full presentation."""
    print("Creating presentation...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    print("  Slide 1: Title")
    add_title_slide(
        prs,
        "Promptable Pathology",
        "Visual Foundation Models for Medical Image Segmentation"
    )
    
    # Slide 2: Problem Statement
    print("  Slide 2: Problem Statement")
    add_content_slide(
        prs,
        "Problem: Manual Annotation is Expensive",
        bullets=[
            "Pathologists spend 15-30 min per slide",
            "Expert annotation costs $50-100 per image",
            "High inter-observer variability",
            "Limited labeled data for rare conditions",
            "",
            "Solution: Visual Foundation Models",
            "Pre-trained on billions of images",
            "Zero-shot capability → no training needed",
            "Promptable interface → flexible interaction"
        ]
    )
    
    # Slide 3: Methods - Pipeline
    print("  Slide 3: Methods")
    pipeline_img = os.path.join(figures_dir, 'academic', 'fig4_method_overview.png')
    add_figure_slide(
        prs,
        "Two-Stage Pipeline: SAM2 → CLIP",
        pipeline_img,
        "Stage 1: SAM2 segments tissue regions | Stage 2: CLIP classifies each region"
    )
    
    # Slide 4: Dataset
    print("  Slide 4: Dataset")
    add_content_slide(
        prs,
        "Dataset: BCSS (Breast Cancer Semantic Segmentation)",
        bullets=[
            "Source: TCGA (The Cancer Genome Atlas)",
            "151 H&E stained images @ 0.25 μm/pixel",
            "5 tissue classes: Tumor, Stroma, Lymphocyte, Necrosis, Blood Vessel",
            "Splits: 85 train / 21 val / 45 test",
            "Challenge: Class imbalance (lymphocytes rare, stroma dominant)"
        ]
    )
    
    # Slide 5: Main Results - Qualitative
    print("  Slide 5: Qualitative Results")
    qual_img = os.path.join(figures_dir, 'qualitative', 'qualitative_method_comparison.png')
    add_figure_slide(
        prs,
        "Segmentation Results: SAM2 vs MedSAM",
        qual_img,
        "4 test images × 6 methods: Original → GT → SAM2 Centroid/Box/Box+Neg → MedSAM"
    )
    
    # Slide 6: Prompt Ablation - Visual
    print("  Slide 6: Prompt Comparison")
    prompt_img = os.path.join(figures_dir, 'qualitative', 'qualitative_prompt_comparison.png')
    add_figure_slide(
        prs,
        "Impact of Prompt Strategy",
        prompt_img,
        "Same image, different prompts → Box+Neg provides best boundary precision"
    )
    
    # Slide 7: Quantitative Segmentation Results
    print("  Slide 7: Segmentation Metrics")
    seg_table = [
        ["Method", "Prompt", "Dice", "IoU", "Δ vs Baseline"],
        ["SAM2", "Centroid", "0.338", "0.236", "baseline"],
        ["SAM2", "Multi-Point", "0.418", "0.287", "+24%"],
        ["SAM2", "Box", "0.553", "0.407", "+64%"],
        ["SAM2", "Box + Neg Points", "0.555", "0.408", "+64% ✓"],
        ["MedSAM", "Box", "0.522", "0.375", "+54%"],
        ["MedSAM", "Box + TTA", "0.536", "0.389", "+59%"],
    ]
    add_results_slide(prs, "Segmentation: Prompt Ablation Results", seg_table, highlight_row=4)
    
    # Slide 8: Per-Class Results
    print("  Slide 8: Per-Class Analysis")
    perclass_img = os.path.join(figures_dir, 'qualitative', 'qualitative_per_class.png')
    add_figure_slide(
        prs,
        "Per-Class Segmentation Performance",
        perclass_img,
        "Necrosis (0.69) easiest → Blood Vessel (0.50) hardest"
    )
    
    # Slide 9: CLIP Results
    print("  Slide 9: CLIP Classification")
    clip_img = os.path.join(figures_dir, 'academic', 'fig2_clip_analysis.png')
    add_figure_slide(
        prs,
        "CLIP Classification: Prompt Engineering",
        clip_img,
        "LLM Few-shot prompts achieve 44.4% accuracy (+91% over baseline)"
    )
    
    # Slide 10: Finetuning Failed
    print("  Slide 10: Finetuning Analysis")
    train_img = os.path.join(figures_dir, 'academic', 'fig3_training_analysis.png')
    add_figure_slide(
        prs,
        "Why Finetuning Failed (-33%)",
        train_img,
        "Small dataset (85 images) → overfitting → catastrophic forgetting"
    )
    
    # Slide 11: Key Findings
    print("  Slide 11: Key Findings")
    findings = [
        ("Zero-shot > Finetuning", "0.555 vs 0.372", "33% performance drop with finetuning on small datasets"),
        ("Box prompts > Point prompts", "+64%", "Spatial context critical for histopathology"),
        ("LLM prompts > Manual prompts", "44.4%", "Few-shot examples improve CLIP accuracy by 91%"),
        ("Text-only > Multimodal", "44% vs 29%", "Images trigger medical jargon CLIP doesn't understand"),
    ]
    add_key_findings_slide(prs, findings)
    
    # Slide 12: Conclusion
    print("  Slide 12: Conclusion")
    add_content_slide(
        prs,
        "Conclusion & Future Work",
        bullets=[
            "Best Configuration:",
            "    SAM2 + Box + Negative Points (0.555 Dice)",
            "    CLIP + LLM Few-shot Prompts (44.4% Accuracy)",
            "",
            "Key Insight: Don't finetune on small medical datasets!",
            "",
            "Future Work:",
            "    → Larger datasets (PathVQA, TCGA-full)",
            "    → Domain-specific VFMs (PathSAM, BiomedCLIP)",
            "    → Interactive refinement loop"
        ]
    )
    
    # Slide 13: Thank You
    print("  Slide 13: Thank You")
    add_title_slide(
        prs,
        "Thank You!",
        "Questions?"
    )
    
    # Save
    output_path = os.path.join(output_dir, 'VFM_Project_Presentation.pptx')
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")
    
    return output_path


if __name__ == '__main__':
    try:
        from pptx import Presentation
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
        from pptx import Presentation
    
    create_presentation()
